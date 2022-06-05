"""
Creates PowerPoint slides with all gifs for a single patient 
on a slide to have a per-patient overview of the predictions.
"""

import os
import sys
sys.path.append(os.path.join('..', '..'))
sys.path.append(os.path.join(__file__, '..', '..', '..'))

from math import ceil
from natsort import natsorted
from pptx import Presentation
from pptx.util import Inches

from utils.config import models_folder


def get_configuration(
    items: int, 
    width: int, 
    height: int, 
    spacing: float, 
    border: float, 
    image_width: int, 
    image_height: int
) -> tuple:
    """
    return number of rows and columns based on the specified height and width
    and the total number of items.
    """
    max_area = 0

    for i in range(1, items+1):
        # get the number of proposed rows and proposed columns
        proposed_columns = i
        proposed_rows = ceil(items/proposed_columns)

        # get the maximum row and column height
        row_height = ((height-border)/proposed_rows)-spacing
        column_width = ((width-border)/proposed_columns)-spacing
        
        # calculate two possible areas with image aspect ratio
        area1 = row_height**2 * image_width/image_height
        area2 = column_width**2 * image_height/image_width

        if area1 < area2:
            area = area1
            column_width = row_height*image_width/image_height
        else:
            area = area2
            row_height = column_width*image_height/image_width
        
        if area > max_area:
            max_area = area
            rows, columns = proposed_rows, proposed_columns
            new_row_height, new_column_width = row_height, column_width

    return rows, columns, new_row_height, new_column_width


if __name__ == '__main__':

    # define model details
    model_subfolder = 'pixel_level'
    model_name = '0001_example_network_0'
    gif_folder = 'clip_prediction_gifs'    # adjust path to folder which contains gifs
    
    # define slide specifics
    width, height = 16, 9 # inches
    spacing = 0.2         # inches
    border = 2*spacing    # inches
    image_width, image_height = 384, 256 # pixels

    # --------------------------------------------------------------------------------------------------

    # create a presentation object
    prs=Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    layout = prs.slide_layouts[6] # empty slide

    # define the directory
    directory = os.path.join(models_folder, model_subfolder, model_name, gif_folder)

    # create a list with gifs in the directory and all different case names
    items = [item for item in natsorted(os.listdir(directory)) if item[-4:] == '.gif']
    cases = natsorted(list(set([item[:8] for item in items])))

    for case in cases:
        # select all gifs for the specific case
        selected_items = [item for item in items if case in item]
        selected_paths = [os.path.join(directory, item) for item in selected_items]
        
        # get the number of rows and columsn, as well as the row and column size (maintaining the aspect ratio)
        rows, columns, row_height, column_width = get_configuration(len(selected_items), width, height, spacing, border, image_width, image_height)

        # add an empty slide
        slide = prs.slides.add_slide(layout)

        i = 0
        for row in range(rows):
            for column in range(columns):
                # get the coordinate position to place gif
                left, top = Inches(border+column*(column_width+spacing)), Inches(border+row*(row_height+spacing))
                # place a gif if there is a position available
                if i < len(selected_items):
                    gif = slide.shapes.add_picture(selected_paths[i], left, top) # add gif
                    # change size of gif
                    gif.width = Inches(column_width)
                    gif.height = Inches(row_height)
                i += 1 
        
        # add the row and column indicators
        for row in range(rows):
            textbox = slide.shapes.add_textbox(Inches(0), Inches(border+(row+0.4)*(row_height+spacing)), Inches(1), Inches(1))
            tf = textbox.text_frame
            tf.text = str(row+1)
        
        for column in range(columns):
            textbox = slide.shapes.add_textbox(Inches(border+(column+0.45)*(column_width+spacing)), Inches(0), Inches(1), Inches(1))
            tf = textbox.text_frame
            tf.text = chr(column+65).upper()
    
        # add the case
        textbox = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(1))
        tf = textbox.text_frame
        tf.text = case

    # save the powerpoint with gifs per patient
    prs.save(os.path.join(directory, "patient_overview.pptx"))